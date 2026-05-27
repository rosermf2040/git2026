from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


SLIDES = [
    {
        "title": "风景如画",
        "subtitle": "一场关于自然之美的旅程",
        "is_title": True,
    },
    {
        "title": "晨曦山岚",
        "body": (
            "天微微亮，薄雾缭绕在群山之间，\n"
            "第一缕阳光穿过云层，洒在松针上，\n"
            "整座山林像被披上了一层金色的轻纱。\n"
            "鸟鸣声此起彼伏，唤醒了沉睡的大地。"
        ),
    },
    {
        "title": "碧海蓝天",
        "body": (
            "海风拂面，咸涩中带着自由的气息，\n"
            "碧蓝的海水与天空连成一线，\n"
            "白色的浪花轻轻拍打着沙滩，\n"
            "海鸥掠过水面，留下一串轻盈的剪影。"
        ),
    },
    {
        "title": "湖光秋色",
        "body": (
            "湖面如镜，倒映着层林尽染的山峦，\n"
            "红的枫、黄的银杏、绿的青松交织成画，\n"
            "微风过处，水波荡漾，落叶轻舞，\n"
            "仿佛走进了一幅流动的水彩画卷。"
        ),
    },
    {
        "title": "雪域星空",
        "body": (
            "夜幕低垂，雪原静谧无声，\n"
            "万千繁星缀满深邃的天幕，\n"
            "银河如绸缎横亘其上，\n"
            "天地之间，唯余呼吸与心跳。"
        ),
    },
    {
        "title": "田园牧歌",
        "body": (
            "金色的麦田随风起伏，\n"
            "远处青瓦白墙的村落升起袅袅炊烟，\n"
            "牛羊在草地上悠闲漫步，\n"
            "这里有最简单也最动人的幸福。"
        ),
    },
    {
        "title": "愿你心中有山海",
        "subtitle": "也愿你眼里有星辰",
        "is_title": True,
    },
]


def build_presentation(path: str) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    bg_color = RGBColor(0x0E, 0x2A, 0x47)
    title_color = RGBColor(0xFF, 0xF5, 0xE1)
    body_color = RGBColor(0xE8, 0xF1, 0xF8)
    accent_color = RGBColor(0xF6, 0xC8, 0x5F)

    blank_layout = prs.slide_layouts[6]

    for slide_data in SLIDES:
        slide = prs.slides.add_slide(blank_layout)

        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = bg_color

        if slide_data.get("is_title"):
            title_box = slide.shapes.add_textbox(
                Inches(1), Inches(2.6), Inches(11.333), Inches(1.6)
            )
            tf = title_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = slide_data["title"]
            run.font.name = "微软雅黑"
            run.font.size = Pt(60)
            run.font.bold = True
            run.font.color.rgb = title_color

            sub_box = slide.shapes.add_textbox(
                Inches(1), Inches(4.4), Inches(11.333), Inches(1)
            )
            sf = sub_box.text_frame
            sf.word_wrap = True
            sp = sf.paragraphs[0]
            sp.alignment = PP_ALIGN.CENTER
            srun = sp.add_run()
            srun.text = slide_data["subtitle"]
            srun.font.name = "微软雅黑"
            srun.font.size = Pt(24)
            srun.font.color.rgb = accent_color
        else:
            title_box = slide.shapes.add_textbox(
                Inches(0.8), Inches(0.7), Inches(11.733), Inches(1.2)
            )
            tf = title_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            run = p.add_run()
            run.text = slide_data["title"]
            run.font.name = "微软雅黑"
            run.font.size = Pt(40)
            run.font.bold = True
            run.font.color.rgb = accent_color

            body_box = slide.shapes.add_textbox(
                Inches(0.8), Inches(2.2), Inches(11.733), Inches(4.6)
            )
            bf = body_box.text_frame
            bf.word_wrap = True
            lines = slide_data["body"].split("\n")
            for idx, line in enumerate(lines):
                bp = bf.paragraphs[0] if idx == 0 else bf.add_paragraph()
                bp.alignment = PP_ALIGN.LEFT
                bp.space_after = Pt(12)
                brun = bp.add_run()
                brun.text = line
                brun.font.name = "微软雅黑"
                brun.font.size = Pt(28)
                brun.font.color.rgb = body_color

    prs.save(path)


if __name__ == "__main__":
    output = "/home/user/git2026/scenic_landscape.pptx"
    build_presentation(output)
    print(f"Generated: {output}")
